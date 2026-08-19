import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen standard dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette (Dark High-Tech Theme)
    BG_DARK = RGBColor(11, 15, 25)       # #0B0F19 Deep Midnight
    CARD_DARK = RGBColor(23, 32, 54)     # #172036 Card Background
    CARD_BORDER = RGBColor(38, 52, 84)   # #263454 Border
    ACCENT_BLUE = RGBColor(59, 130, 246) # #3B82F6 Vibrant Blue
    ACCENT_CYAN = RGBColor(6, 182, 212)  # #06B6D4 Cyan
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981 Emerald Green
    ACCENT_ORANGE = RGBColor(245, 158, 11)# #F59E0B Amber
    TEXT_WHITE = RGBColor(255, 255, 255) # Pure White
    TEXT_LIGHT = RGBColor(226, 232, 240) # Slate 200
    TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400

    def add_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background() # No border
        return bg

    def add_header(slide, title_text, category_text="MYMOBPAY • COLLEGE PROJECT PRESENTATION"):
        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(9.5)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_CYAN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.65))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, title=None, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_DARK
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = "Arial"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE

        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide1)

    # Accent Top Bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(1.5), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_CYAN
    bar.line.fill.background()

    # Project Title
    tbox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "MyMobPay"
    p1.font.name = "Arial"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Zero-Commission Direct P2P UPI Payment Gateway & SaaS Infrastructure"
    p2.font.name = "Arial"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_BLUE

    p3 = tf.add_paragraph()
    p3.text = "An Automated Bank SMS-Reconciliation Architecture Eliminating Aggregator Fees (0% MDR)"
    p3.font.name = "Arial"
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED

    # Highlights 3 Cards
    cards_data = [
        ("0% Transaction Fee", "Direct P2P routing from customer to merchant account without 2-3% aggregator cut.", ACCENT_GREEN),
        ("Sub-2s Reconciliation", "Automated SMS-to-Webhook parsing verifying payments instantly with zero manual effort.", ACCENT_CYAN),
        ("Full-Stack Modern SaaS", "Next.js 14 App Router, Supabase PostgreSQL, Tailwind CSS, and Cloudflare Edge.", ACCENT_BLUE)
    ]
    for i, (ctitle, cdesc, ccolor) in enumerate(cards_data):
        c_left = Inches(0.8 + i * 4.0)
        c_card = add_card(slide1, c_left, Inches(3.6), Inches(3.7), Inches(2.1), border_color=ccolor)
        
        tb = slide1.shapes.add_textbox(c_left + Inches(0.25), Inches(3.85), Inches(3.2), Inches(1.6))
        tf_c = tb.text_frame
        tf_c.word_wrap = True
        p_ct = tf_c.paragraphs[0]
        p_ct.text = ctitle
        p_ct.font.name = "Arial"
        p_ct.font.size = Pt(14)
        p_ct.font.bold = True
        p_ct.font.color.rgb = ccolor
        
        p_cd = tf_c.add_paragraph()
        p_cd.text = "\n" + cdesc
        p_cd.font.name = "Arial"
        p_cd.font.size = Pt(11)
        p_cd.font.color.rgb = TEXT_LIGHT

    # Metadata Footer Box
    meta_box = slide1.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9))
    tf_m = meta_box.text_frame
    tf_m.word_wrap = True
    pm1 = tf_m.paragraphs[0]
    pm1.text = "Submitted by: Final Year Engineering Project Team  |  Domain: FinTech & Full-Stack Cloud Computing"
    pm1.font.name = "Arial"
    pm1.font.size = Pt(11)
    pm1.font.bold = True
    pm1.font.color.rgb = TEXT_WHITE
    pm2 = tf_m.add_paragraph()
    pm2.text = "Live Production URL: https://mymob.tech  |  Tech: Next.js 14, Supabase (PostgreSQL), RESTful Webhooks"
    pm2.font.name = "Arial"
    pm2.font.size = Pt(10)
    pm2.font.color.rgb = ACCENT_CYAN

    # ==========================================
    # SLIDE 2: INTRODUCTION & BACKGROUND
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide2)
    add_header(slide2, "Introduction: The UPI Ecosystem & Market Context", "01. PROJECT OVERVIEW")

    add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "The Rise of Unified Payments Interface (UPI)")
    tb2_1 = slide2.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(5.1), Inches(4.5))
    tf2_1 = tb2_1.text_frame
    tf2_1.word_wrap = True
    bullets_2_1 = [
        "Revolutionary Growth: UPI processes over 14+ Billion transactions monthly in India, becoming the world's largest real-time retail payment system.",
        "Zero Cost for P2P: Direct person-to-person UPI transactions incur ₹0 fees from NPCI.",
        "Merchant Dilemma: Traditional payment gateways (Razorpay, Cashfree, PayU) treat UPI as a commercial gateway transaction and charge 1.8% to 3.0% per transaction + 18% GST.",
        "Margin Destruction: For thin-margin B2B startups, creators, digital product sellers, and retail stores, 2-3% revenue loss directly wipes out net profit."
    ]
    for i, b in enumerate(bullets_2_1):
        p = tf2_1.paragraphs[0] if i == 0 else tf2_1.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "Project Objectives & Scope")
    tb2_2 = slide2.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.2), Inches(4.5))
    tf2_2 = tb2_2.text_frame
    tf2_2.word_wrap = True
    bullets_2_2 = [
        "Democratize B2B Payments: Build a self-hosted or SaaS-ready payment infrastructure with 0% gateway commission.",
        "Direct Bank-to-Bank Settlement: Direct settlement to the merchant's UPI VPA instantly without T+1/T+2 holding days.",
        "Automated Reconciliation: Solve the fundamental flaw of static QR codes (manual payment proof checking) using an automated SMS-parsing Webhook engine.",
        "Enterprise-Grade UX: Provide dynamic QR generation, mobile UPI app deep-linking (GPay/PhonePe/Paytm), real-time status polling, and comprehensive analytics."
    ]
    for i, b in enumerate(bullets_2_2):
        p = tf2_2.paragraphs[0] if i == 0 else tf2_2.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: PROBLEM STATEMENT
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide3)
    add_header(slide3, "Problem Statement: Pain Points of Current Solutions", "02. NEED FOR THE PROJECT")

    problems = [
        ("High Aggregator Commissions (MDR)", "Traditional gateways charge 2% - 3% + GST on every transaction. On ₹10,00,000 monthly turnover, merchants lose ₹25,000 - ₹35,000 pure profit to middlemen.", ACCENT_ORANGE),
        ("Settlement Delays & Fund Freezes", "Aggregators hold merchant funds for T+1 or T+2 working days. High-risk or sudden-volume merchants face arbitrary account freezes and withheld payouts.", ACCENT_ORANGE),
        ("Manual Static QR Verification Hell", "Static paper QR codes require staff to manually inspect customer phone screenshots, creating long checkout queues, fake screenshot fraud, and accounting chaos.", ACCENT_ORANGE),
        ("Complex KYC & High Entry Barriers", "Small businesses and individual founders struggle with enterprise KYC, onboarding delays, recurring setup fees, and restrictive business categories.", ACCENT_ORANGE)
    ]

    for i, (title, desc, color) in enumerate(problems):
        row = i // 2
        col = i % 2
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 5.9)
        add_card(slide3, left, top, Inches(5.7), Inches(2.4), title, color)
        
        tb = slide3.shapes.add_textbox(left + Inches(0.25), top + Inches(0.7), Inches(5.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 4: THE PROPOSED SOLUTION (MYMOBPAY)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide4)
    add_header(slide4, "The Solution: MyMobPay Architecture & Value Proposition", "03. PROPOSED SYSTEM")

    sol_cards = [
        ("1. Direct P2P Payment Routing", "Generates dynamic UPI URLs linking directly to merchant's bank VPA. 100% of customer funds transfer immediately without escrow.", ACCENT_GREEN),
        ("2. Automated SMS Verification Engine", "Secure Android background daemon forwards incoming bank credit SMS to backend webhook. Regex engine parses amount & UTR in real-time.", ACCENT_CYAN),
        ("3. Sub-2-Second Reconciliation", "State machine updates order status from 'created' -> 'pending' -> 'verified' within milliseconds of bank credit alert.", ACCENT_BLUE),
        ("4. Seamless Cross-Platform UX", "Desktop shows interactive dynamic QR code; Mobile switches to native one-click intent launcher for Google Pay, PhonePe, Paytm, BHIM.", ACCENT_GREEN),
        ("5. Multi-Tenant Merchant SaaS", "Merchants manage custom branding, staff role-based access control, analytics, webhook endpoints, and API keys through a centralized dashboard.", ACCENT_CYAN),
        ("6. Zero Intermediary Escrow", "No funds ever touch MyMobPay servers. Eliminates regulatory compliance bottlenecks, chargeback holds, and gateway escrow risks.", ACCENT_BLUE)
    ]

    for i, (title, desc, color) in enumerate(sol_cards):
        row = i // 3
        col = i % 3
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 3.9)
        add_card(slide4, left, top, Inches(3.7), Inches(2.4), title, color)
        
        tb = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), Inches(3.3), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 5: SYSTEM ARCHITECTURE & DATA FLOW
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide5)
    add_header(slide5, "System Architecture: End-to-End Workflow & Modules", "04. SYSTEM ARCHITECTURE")

    steps = [
        ("Step 1: Order Initiation", "Client or external merchant site calls `/api/orders`.\nGenerates unique Order ID & custom UPI payment intent string (`upi://pay?pa=...`).", ACCENT_BLUE),
        ("Step 2: Customer Checkout", "Desktop: High-fidelity QR rendered.\nMobile: Native intent links for GPay / PhonePe.\nCustomer authorizes payment via UPI PIN in bank app.", ACCENT_CYAN),
        ("Step 3: Instant Bank Credit", "Funds are instantly credited from Customer Bank Account to Merchant Bank Account via NPCI UPI network.", ACCENT_GREEN),
        ("Step 4: SMS Forwarder Webhook", "Merchant's Android device receives bank SMS.\nForwarder app triggers secure POST request to `/api/webhook/sms` with HMAC verification.", ACCENT_ORANGE),
        ("Step 5: Parsing & Reconciliation", "Backend Regex engine identifies bank, parses amount & UTR, queries pending orders, transitions status to 'verified', logs timestamp.", ACCENT_CYAN),
        ("Step 6: Real-Time Client Update", "Customer `/status/[orderId]` UI updates immediately via smart polling / Supabase Realtime.\nDispatches HMAC-signed merchant webhook.", ACCENT_GREEN)
    ]

    for i, (title, desc, color) in enumerate(steps):
        row = i // 3
        col = i % 3
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 3.9)
        add_card(slide5, left, top, Inches(3.7), Inches(2.4), title, color)
        
        tb = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), Inches(3.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 6: THE SECRET SAUCE - REGEX PARSING ENGINE
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide6)
    add_header(slide6, "The Core Engine: Bank SMS Regex Parser (`parseSms.js`)", "05. TECHNICAL DEEP DIVE")

    add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "Supported Indian Banking Gateways")
    tb6_1 = slide6.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(5.1), Inches(4.5))
    tf6_1 = tb6_1.text_frame
    tf6_1.word_wrap = True
    banks = [
        "State Bank of India (SBI): Detects 'credited by Rs.', UPI Ref numbers, account masks.",
        "HDFC Bank: Extracts 'deposited', 'credited to account', 'UPI/Ref/...' structures.",
        "ICICI Bank: Parses 'credited with INR', 12-digit UTR sequences, timestamp logs.",
        "Axis Bank: Matches 'received in A/c', UPI transaction references.",
        "Indian Overseas Bank (IOB) & Generic Fallback: Robust heuristic regex extracting standard currency formats `(?:Rs\\.?|INR)\\s*([\\d,]+\\.?\\d*)`."
    ]
    for i, b in enumerate(banks):
        p = tf6_1.paragraphs[0] if i == 0 else tf6_1.add_paragraph()
        p.text = "✓ " + b
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(8)

    add_card(slide6, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "SMS Parsing & Fraud Guard Code Architecture")
    tb6_2 = slide6.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.2), Inches(4.5))
    tf6_2 = tb6_2.text_frame
    tf6_2.word_wrap = True
    code_expl = [
        "Stateless Token Security: Requests require strict `webhookSecret` verification in payload headers.",
        "Anti-Collision Matching: Matches against `orders` where `status = 'pending'` and `amount = parsed_amount` ordered by `created_at DESC` within a 15-minute validity window.",
        "Unique UTR Enforcement: Extracted 12-digit UPI UTR number is stored under PostgreSQL `UNIQUE` constraint preventing duplicate SMS replay attacks.",
        "Atomic State Transitions: Database mutations run with service-role privileges bypassing client-side tampering."
    ]
    for i, b in enumerate(code_expl):
        p = tf6_2.paragraphs[0] if i == 0 else tf6_2.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 7: SECURITY & FRAUD PREVENTION MECHANICS
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide7)
    add_header(slide7, "Security Architecture: Anti-Fraud & Data Protection", "06. SECURITY & COMPLIANCE")

    sec_cards = [
        ("Duplicate UTR Replay Protection", "Every bank SMS contains a unique 12-digit UTR / Reference ID. Supabase enforces a database-level UNIQUE constraint on the `utr` column. Duplicate requests return `DUPLICATE_UTR` error and are discarded.", ACCENT_GREEN),
        ("HMAC Webhook Authentication", "Incoming SMS forwarder requests must include the pre-shared `webhookSecret` in headers/payload, shielding the endpoint from rogue network scans and spam.", ACCENT_CYAN),
        ("PCI-DSS Compliance by Design", "No customer debit card numbers, CVVs, UPI PINs, or bank passwords ever enter or touch MyMobPay servers. The system only handles public VPAs and standard NPCI intent strings.", ACCENT_BLUE),
        ("Role-Based Access Control (RBAC)", "Multi-tenant merchant isolation with secure authentication via Supabase Auth. Staff permissions hierarchy (Admin vs Cashier vs Developer).", ACCENT_GREEN)
    ]

    for i, (title, desc, color) in enumerate(sec_cards):
        row = i // 2
        col = i % 2
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 5.9)
        add_card(slide7, left, top, Inches(5.7), Inches(2.4), title, color)
        
        tb = slide7.shapes.add_textbox(left + Inches(0.25), top + Inches(0.7), Inches(5.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 8: DATABASE SCHEMA & ENTITY MODEL
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide8)
    add_header(slide8, "Database Design: PostgreSQL Relational Schema", "07. DATA MODELING")

    schema_tables = [
        ("Table: `merchants`", "• `id`: UUID (Primary Key / Auth ID)\n• `business_name`: Text\n• `upi_id`: Text (Receiving VPA)\n• `phone_number`: Text (10 digits)\n• `subscription_status`: Enum ('active', 'trial')\n• `sandbox_mode`: Boolean (Simulation flag)\n• `theme_color`: Hex Code (#3B82F6)\n• `created_at`: Timestamp"),
        ("Table: `orders`", "• `id`: Text (Short unique ID, e.g. O1B2)\n• `merchant_id`: UUID (Foreign Key)\n• `amount`: Numeric(10,2) (Transaction value)\n• `status`: Enum ('pending','verified','rejected')\n• `utr`: Text (UNIQUE constraint)\n• `customer_name`: Text\n• `customer_phone`: Text\n• `verified_at`: Timestamp"),
        ("Table: `admin_settings` & `staff`", "• `id`: UUID / Text Key\n• `user_id`: UUID (Foreign Key to auth.users)\n• `role`: Enum ('super_admin','staff','cashier')\n• `webhook_url`: Text (Developer Notification)\n• `webhook_secret`: Text (HMAC Key)\n• `totp_secret`: Text (Two-factor auth)\n• `is_active`: Boolean")
    ]

    for i, (tbl_name, tbl_fields) in enumerate(schema_tables):
        left = Inches(0.8 + i * 4.0)
        add_card(slide8, left, Inches(1.5), Inches(3.7), Inches(5.3), tbl_name, ACCENT_BLUE)
        
        tb = slide8.shapes.add_textbox(left + Inches(0.25), Inches(2.1), Inches(3.2), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = tbl_fields
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 9: FULL TECH STACK & TOOLS
    # ==========================================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide9)
    add_header(slide9, "Technology Stack & Development Environment", "08. TECH STACK")

    tech_cards = [
        ("Frontend Technologies", "• Next.js 14 (App Router & Server Components)\n• Tailwind CSS (Utility-First Responsive UI)\n• Lucide Icons (Modern Iconography)\n• React-QR-Code (Dynamic Vector QR Rendering)\n• Canvas Confetti (Celebratory Success UX)", ACCENT_CYAN),
        ("Backend & Edge Computing", "• Next.js Edge & Serverless API Handlers\n• Node.js Crypto (HMAC SHA-256 Signature Tokenization)\n• Cloudflare Workers (High-Throughput Forwarding)\n• RESTful JSON API Protocols", ACCENT_BLUE),
        ("Database & Cloud Infrastructure", "• Supabase (Managed PostgreSQL DB)\n• Row-Level Security (RLS) Policy Layer\n• Supabase Service-Role Admin Client\n• Vercel Global Edge Network Deployment", ACCENT_GREEN),
        ("Android Daemon Forwarder", "• Background SMS Listener Service\n• BroadcastReceiver on `SMS_RECEIVED`\n• Auto-retry queue for poor connectivity\n• Secure HTTPS POST payload dispatch", ACCENT_ORANGE)
    ]

    for i, (title, desc, color) in enumerate(tech_cards):
        row = i // 2
        col = i % 2
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 5.9)
        add_card(slide9, left, top, Inches(5.7), Inches(2.4), title, color)
        
        tb = slide9.shapes.add_textbox(left + Inches(0.25), top + Inches(0.7), Inches(5.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(10.5)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 10: USER EXPERIENCE & DUAL-SCREEN INTERFACE
    # ==========================================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide10)
    add_header(slide10, "UI/UX Design: Intelligent Dual-Device Checkout", "09. USER EXPERIENCE")

    add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.3), "🖥️ Desktop Checkout Interface")
    tb10_1 = slide10.shapes.add_textbox(Inches(1.05), Inches(2.1), Inches(5.1), Inches(4.5))
    tf10_1 = tb10_1.text_frame
    tf10_1.word_wrap = True
    desktop_features = [
        "Real-Time Dynamic QR Code: Encodes NPCI URI parameters (`pa`, `pn`, `am`, `cu`, `tn`, `tr`).",
        "Visual Laser Scanner Animation: Polished fintech aesthetics instilling trust and credibility.",
        "Auto-Expiry Countdown Timer: 10-minute validity preventing stale pending order buildup.",
        "Instant Verification Transition: Smooth animation into celebratory verified badge with transaction summary receipt."
    ]
    for i, b in enumerate(desktop_features):
        p = tf10_1.paragraphs[0] if i == 0 else tf10_1.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    add_card(slide10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.3), "📱 Mobile Checkout Experience")
    tb10_2 = slide10.shapes.add_textbox(Inches(7.05), Inches(2.1), Inches(5.2), Inches(4.5))
    tf10_2 = tb10_2.text_frame
    tf10_2.word_wrap = True
    mobile_features = [
        "1-Click App Deep Linking: Direct trigger buttons for Google Pay, PhonePe, Paytm, BHIM, and Cred.",
        "Zero Typing Friction: Pre-fills receiving VPA and exact payable amount in UPI app directly.",
        "Smart Background App Switching: Customer returns to browser tab where status screen automatically detects settlement.",
        "Direct Direct-Signup Integration: Streamlined merchant onboarding with zero OTP friction."
    ]
    for i, b in enumerate(mobile_features):
        p = tf10_2.paragraphs[0] if i == 0 else tf10_2.add_paragraph()
        p.text = "• " + b
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 11: MERCHANT DASHBOARD & ADMIN CAPABILITIES
    # ==========================================
    slide11 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide11)
    add_header(slide11, "Merchant Console & Super Admin Capabilities", "10. MANAGEMENT CONSOLE")

    m_features = [
        ("Live Revenue & Order Analytics", "Real-time KPI metrics: Total Gross Volume (₹), Today's Collections, Successful vs Pending Orders, and Conversion Rate graphs.", ACCENT_GREEN),
        ("Manual Override & Dispute Resolution", "Full transaction audit ledger with search, filtering by date/status, and manual 'Force Verify' / 'Reject' administrator overrides.", ACCENT_BLUE),
        ("Sandbox Simulation Suite", "Built-in mock payment simulator allowing developers and staff to test end-to-end webhook flows without spending actual money.", ACCENT_CYAN),
        ("API Keys & Webhook Management", "Generate HMAC-signed developer API keys, configure outgoing webhook listener URLs, and inspect request delivery logs.", ACCENT_ORANGE)
    ]

    for i, (title, desc, color) in enumerate(m_features):
        row = i // 2
        col = i % 2
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 5.9)
        add_card(slide11, left, top, Inches(5.7), Inches(2.4), title, color)
        
        tb = slide11.shapes.add_textbox(left + Inches(0.25), top + Inches(0.7), Inches(5.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 12: COMPARATIVE ANALYSIS (MYMOBPAY VS TRADITIONAL)
    # ==========================================
    slide12 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide12)
    add_header(slide12, "Comparative Analysis: MyMobPay vs Traditional Gateways", "11. COMPETITIVE ADVANTAGE")

    # Table comparison
    rows = [
        ("Feature / Metric", "Traditional Gateways (Razorpay/Cashfree)", "MyMobPay (Proposed Solution)"),
        ("Transaction Fee (MDR)", "2.0% - 3.0% + 18% GST per transaction", "0% (Zero Gateway Commission)"),
        ("Settlement Speed", "T+1 or T+2 Business Days (Delayed)", "Instant (0 Seconds directly in bank)"),
        ("Intermediary Custody", "Funds held in Aggregator Escrow", "Direct P2P (Zero Custodial Risk)"),
        ("Verification Method", "Aggregator Webhook Callback", "Automated Bank SMS Parsing (< 2s)"),
        ("Account Freezes", "Frequent due to automated risk triggers", "Impossible (Bank retains full custody)"),
        ("Monthly Cost Model", "High variable commission", "Predictable flat SaaS subscription")
    ]

    # Create visual table
    t_top = Inches(1.5)
    t_left = Inches(0.8)
    t_width = Inches(11.7)
    t_height = Inches(5.2)

    table_shape = slide12.shapes.add_table(len(rows), 3, t_left, t_top, t_width, t_height)
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(4.0)

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            if r_idx == 0:
                p.font.bold = True
                p.font.size = Pt(11.5)
                p.font.color.rgb = TEXT_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            else:
                p.font.size = Pt(10.5)
                if c_idx == 2:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_CYAN
                elif c_idx == 1:
                    p.font.color.rgb = RGBColor(248, 113, 113) # Light red
                else:
                    p.font.color.rgb = TEXT_LIGHT
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_DARK if r_idx % 2 == 0 else RGBColor(18, 25, 43)

    # ==========================================
    # SLIDE 13: FUTURE ENHANCEMENTS & ROADMAP
    # ==========================================
    slide13 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide13)
    add_header(slide13, "Future Scope, Commercialization & Technical Roadmap", "12. FUTURE WORK")

    roadmap_items = [
        ("AI-Powered Anomaly & Fraud Detection", "Deploy machine learning models to detect irregular transaction bursts, spoofed SMS patterns, and unusual geographic latency.", ACCENT_BLUE),
        ("Hardware Soundbox Integration", "Integrate low-cost IoT Bluetooth / 4G Voice Soundboxes for instant audio payment announcements at physical retail counters.", ACCENT_CYAN),
        ("Multi-Bank Dynamic VPA Failover", "Automatically distribute high-frequency incoming transactions across multiple merchant bank accounts to avoid daily UPI velocity limits.", ACCENT_GREEN),
        ("Pre-built E-Commerce Plugins", "Release official 1-click plugins for WooCommerce (WordPress), Shopify, Magento, and WHMCS.", ACCENT_ORANGE)
    ]

    for i, (title, desc, color) in enumerate(roadmap_items):
        row = i // 2
        col = i % 2
        top = Inches(1.5 + row * 2.7)
        left = Inches(0.8 + col * 5.9)
        add_card(slide13, left, top, Inches(5.7), Inches(2.4), title, color)
        
        tb = slide13.shapes.add_textbox(left + Inches(0.25), top + Inches(0.7), Inches(5.2), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_LIGHT

    # ==========================================
    # SLIDE 14: CONCLUSION & Q&A
    # ==========================================
    slide14 = prs.slides.add_slide(blank_slide_layout)
    add_slide_background(slide14)

    # Center card for conclusion
    c_card = add_card(slide14, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.3), border_color=ACCENT_BLUE)
    
    tb14 = slide14.shapes.add_textbox(Inches(1.8), Inches(1.5), Inches(9.7), Inches(4.7))
    tf14 = tb14.text_frame
    tf14.word_wrap = True
    
    p = tf14.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "CONCLUSION & SUMMARY"
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN

    bullets_c = [
        "MyMobPay demonstrates that modern full-stack cloud computing (Next.js 14 + Supabase) combined with mobile daemon webhooks can completely bypass traditional payment aggregator monopolies.",
        "Delivers 100% genuine 0% MDR transaction routing with real-time automated bank reconciliation in under 2 seconds.",
        "Eliminates custodial risk, chargeback delays, and high fees for thousands of small businesses, creators, and developers.",
        "Fully deployed, production-tested, and live at: https://mymob.tech"
    ]
    
    for b in bullets_c:
        p_b = tf14.add_paragraph()
        p_b.text = "\n✓ " + b
        p_b.font.name = "Arial"
        p_b.font.size = Pt(12.5)
        p_b.font.color.rgb = TEXT_LIGHT

    p_qa = tf14.add_paragraph()
    p_qa.alignment = PP_ALIGN.CENTER
    p_qa.text = "\nThank You! Questions & Discussion Welcome."
    p_qa.font.name = "Arial"
    p_qa.font.size = Pt(18)
    p_qa.font.bold = True
    p_qa.font.color.rgb = ACCENT_GREEN

    # Save presentation
    output_path = os.path.join(os.getcwd(), "MyMobPay_College_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
